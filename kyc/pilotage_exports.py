"""
kyc/pilotage_exports.py
Génération des rapports PDF et PPTX pour le Pilotage KYC.
Utilise ReportLab (PDF) et python-pptx (PPTX).
"""

import io
import os
import math
from datetime import datetime

from django.conf import settings
from django.http import HttpResponse
from django.utils import timezone

                                                                                
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm, mm
from reportlab.platypus import (
    BaseDocTemplate, Frame, PageTemplate, Paragraph, Spacer, Table,
    TableStyle, KeepTogether, HRFlowable,
)
from reportlab.platypus.flowables import Flowable
from reportlab.graphics.shapes import Drawing, Rect, String, Line
from reportlab.graphics import renderPDF
from reportlab.graphics.charts.barcharts import HorizontalBarChart
from reportlab.graphics.charts.piecharts import Pie
from reportlab.pdfgen import canvas as pdfcanvas
from reportlab.lib.utils import ImageReader

                                                                                
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

                                                                                
BOA_GREEN       = colors.HexColor("#00965E")                      
BOA_GREEN_LIGHT = colors.HexColor("#E8F5E9")                    
BOA_GREEN_MED   = colors.HexColor("#81C784")
BOA_DARK        = colors.HexColor("#0f172a")
BOA_BLUE        = colors.HexColor("#1B2A4A")                    
BOA_BLUE_LIGHT  = colors.HexColor("#2563eb")                
BOA_SLATE       = colors.HexColor("#64748b")
BOA_RED         = colors.HexColor("#ef4444")
BOA_RED_LIGHT   = colors.HexColor("#fee2e2")
BOA_AMBER       = colors.HexColor("#f59e0b")
BOA_AMBER_LIGHT = colors.HexColor("#fef3c7")
BOA_WHITE       = colors.white
BOA_GRAY        = colors.HexColor("#f1f5f9")
BOA_BORDER      = colors.HexColor("#e2e8f0")

                                                                               
PPTX_GREEN       = RGBColor(0x00, 0x96, 0x5E)                            
PPTX_GREEN_DEEP  = RGBColor(0x00, 0x6B, 0x42)                            
PPTX_GREEN_SOFT  = RGBColor(0xE8, 0xF5, 0xEE)                            
PPTX_DARK        = RGBColor(0x1B, 0x2A, 0x4A)                            
PPTX_DARK_DEEP   = RGBColor(0x0F, 0x17, 0x2A)                            
PPTX_RED         = RGBColor(0xEF, 0x44, 0x44)                 
PPTX_RED_SOFT    = RGBColor(0xFF, 0xEB, 0xEB)                    
PPTX_AMBER       = RGBColor(0xF5, 0x9E, 0x0B)                    
PPTX_AMBER_SOFT  = RGBColor(0xFE, 0xF3, 0xC7)                    
PPTX_SLATE       = RGBColor(0x64, 0x74, 0x8B)                 
PPTX_GRAY        = RGBColor(0xE2, 0xE8, 0xF0)                      
PPTX_GRAY_BG     = RGBColor(0xF8, 0xFA, 0xFC)                    
PPTX_GRAY_TEXT   = RGBColor(0x94, 0xA3, 0xB8)                     
PPTX_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
PPTX_GREEN_LIGHT = RGBColor(0xBB, 0xF7, 0xD0)                      
PPTX_GREEN_PALE  = RGBColor(0xDC, 0xFC, 0xE7)                        
PPTX_GOLD        = RGBColor(0xF5, 0xC5, 0x18)                      

LOGO_PATH = os.path.join(settings.MEDIA_ROOT, "images", "boa_logo.png")


def _logo_path():
    if os.path.isfile(LOGO_PATH):
        return LOGO_PATH
    return None


def _fmt_pct(rate, dash="—"):
    """Taux en % : sans décimale s'il est entier (cas complétude, tronquée),
    sinon une décimale."""
    if rate is None:
        return dash
    rate = float(rate)
    return f"{rate:.0f}%" if rate.is_integer() else f"{rate:.1f}%"


def _rate_color(rate, threshold):
    """Retourne la couleur en fonction du taux (vert ou orange uniquement)."""
    if rate is None:
        return BOA_SLATE
    return BOA_AMBER if rate < threshold else BOA_GREEN


def _rate_color_pptx(rate, threshold):
    if rate is None:
        return PPTX_SLATE
    return PPTX_RED if rate < threshold else PPTX_GREEN


                                                                                 
                                                 
                                                                                 

class RateBar(Flowable):
    """Petite barre de progression colorée pour les tableaux."""
    def __init__(self, rate, threshold, width=60, height=8):
        super().__init__()
        self.rate = rate
        self.threshold = threshold
        self.width = width
        self.height = height

    def draw(self):
        c = self.canv
        rate = self.rate if self.rate is not None else 0
        fill_pct = min(max(rate, 0), 100) / 100
        bar_color = _rate_color(self.rate, self.threshold)

                   
        c.setFillColor(BOA_GRAY)
        c.roundRect(0, 0, self.width, self.height, 3, fill=1, stroke=0)
                       
        c.setFillColor(bar_color)
        c.roundRect(0, 0, self.width * fill_pct, self.height, 3, fill=1, stroke=0)

    def wrap(self, available_width, available_height):
        return self.width, self.height


class SetSectionTitle(Flowable):
    """Flowable permettant de changer le titre de la section courante sur le canvas."""
    def __init__(self, title):
        super().__init__()
        self.title = title

    def draw(self):
        self.canv._current_section_title = self.title

    def wrap(self, availWidth, availHeight):
        return 0, 0


                                                                                 
                 
                                                                                 

def _build_pdf_styles():
    base = getSampleStyleSheet()
    styles = {
        "title": ParagraphStyle(
            "title", parent=base["Normal"],
            fontName="Helvetica-Bold", fontSize=22, textColor=BOA_WHITE,
            leading=26, alignment=TA_LEFT,
        ),
        "subtitle": ParagraphStyle(
            "subtitle", parent=base["Normal"],
            fontName="Helvetica", fontSize=10, textColor=BOA_GREEN_MED,
            leading=14, alignment=TA_LEFT,
        ),
        "section_title": ParagraphStyle(
            "section_title", parent=base["Normal"],
            fontName="Helvetica-Bold", fontSize=12, textColor=BOA_WHITE,
            leading=16, alignment=TA_LEFT,
        ),
        "kpi_label": ParagraphStyle(
            "kpi_label", parent=base["Normal"],
            fontName="Helvetica-Bold", fontSize=7, textColor=BOA_SLATE,
            leading=10, alignment=TA_CENTER, spaceAfter=2,
        ),
        "kpi_value": ParagraphStyle(
            "kpi_value", parent=base["Normal"],
            fontName="Helvetica-Bold", fontSize=22, textColor=BOA_DARK,
            leading=26, alignment=TA_CENTER,
        ),
        "kpi_value_red": ParagraphStyle(
            "kpi_value_red", parent=base["Normal"],
            fontName="Helvetica-Bold", fontSize=22, textColor=BOA_AMBER,
            leading=26, alignment=TA_CENTER,
        ),
        "th": ParagraphStyle(
            "th", parent=base["Normal"],
            fontName="Helvetica-Bold", fontSize=7, textColor=BOA_WHITE,
            leading=10, alignment=TA_LEFT,
        ),
        "td": ParagraphStyle(
            "td", parent=base["Normal"],
            fontName="Helvetica", fontSize=7, textColor=BOA_DARK,
            leading=10, alignment=TA_LEFT,
        ),
        "td_right": ParagraphStyle(
            "td_right", parent=base["Normal"],
            fontName="Helvetica", fontSize=7, textColor=BOA_DARK,
            leading=10, alignment=TA_RIGHT,
        ),
        "td_bold": ParagraphStyle(
            "td_bold", parent=base["Normal"],
            fontName="Helvetica-Bold", fontSize=7, textColor=BOA_DARK,
            leading=10, alignment=TA_LEFT,
        ),
        "rate_tag": ParagraphStyle(
            "rate_tag", parent=base["Normal"],
            fontName="Helvetica-Bold", fontSize=8, textColor=BOA_WHITE,
            leading=11, alignment=TA_CENTER,
        ),
        "footer": ParagraphStyle(
            "footer", parent=base["Normal"],
            fontName="Helvetica", fontSize=7, textColor=BOA_SLATE,
            leading=10, alignment=TA_CENTER,
        ),
        "body": ParagraphStyle(
            "body", parent=base["Normal"],
            fontName="Helvetica", fontSize=9, textColor=BOA_DARK,
            leading=13, alignment=TA_LEFT,
        ),
        "chart_title": ParagraphStyle(
            "chart_title", parent=base["Normal"],
            fontName="Helvetica-Bold", fontSize=9, textColor=BOA_DARK,
            leading=13, alignment=TA_LEFT, spaceAfter=6,
        ),
        "toc_title": ParagraphStyle(
            "toc_title", parent=base["Normal"],
            fontName="Helvetica-Bold", fontSize=16, textColor=BOA_BLUE,
            leading=22, alignment=TA_LEFT, spaceAfter=12,
        ),
        "toc_entry": ParagraphStyle(
            "toc_entry", parent=base["Normal"],
            fontName="Helvetica-Bold", fontSize=10, textColor=BOA_DARK,
            leading=16, alignment=TA_LEFT, spaceBefore=4,
        ),
        "toc_sub": ParagraphStyle(
            "toc_sub", parent=base["Normal"],
            fontName="Helvetica", fontSize=9, textColor=BOA_SLATE,
            leading=14, alignment=TA_LEFT, leftIndent=20,
        ),
        "page_tag": ParagraphStyle(
            "page_tag", parent=base["Normal"],
            fontName="Helvetica-Bold", fontSize=8, textColor=BOA_WHITE,
            leading=11, alignment=TA_RIGHT,
        ),
    }
    return styles


                                                                                 
                                            
                                                                                 

def _build_completeness_bar_chart(rows, threshold, max_items=15):
    """Construit un HorizontalBarChart ReportLab pour la complétude."""
    if not rows:
        return None

    items = rows[:max_items]
    labels = [f"{r['type']} - {r['field_label'][:22]}" for r in items]
    values = [r["rate"] if r["rate"] is not None else 0 for r in items]
    n = len(items)

    chart_height = max(120, n * 18 + 40)
    chart_width = 480

    drawing = Drawing(chart_width, chart_height)
    bc = HorizontalBarChart()
    bc.x = 160
    bc.y = 20
    bc.width = 260
    bc.height = chart_height - 40
    bc.data = [values]
    bc.categoryAxis.categoryNames = labels
    bc.categoryAxis.labels.fontSize = 7
    bc.categoryAxis.labels.fontName = "Helvetica"
    bc.categoryAxis.labels.dx = -4
    bc.categoryAxis.labels.textAnchor = "end"
    bc.valueAxis.valueMin = 0
    bc.valueAxis.valueMax = 100
    bc.valueAxis.valueStep = 20
    bc.valueAxis.labels.fontSize = 7
    bc.valueAxis.labels.fontName = "Helvetica"
    bc.bars[0].strokeColor = None

                                          
    bc.barLabelFormat = '%d%%'
    bc.barLabels.fontName = "Helvetica-Bold"
    bc.barLabels.fontSize = 6
    bc.barLabels.fillColor = BOA_BLUE
    bc.barLabels.boxAnchor = 'w'
    bc.barLabels.textAnchor = 'start'
    bc.barLabels.dx = 4
    bc.barLabels.nudge = 4

                                                          
    for i, rate in enumerate(values):
        bc.bars[0, i].fillColor = _rate_color(rate, threshold)

    drawing.add(bc)

                              
    seuil_x = bc.x + (threshold / 100) * bc.width
    line = Line(seuil_x, bc.y, seuil_x, bc.y + bc.height,
                strokeColor=BOA_RED, strokeWidth=1, strokeDashArray=[3, 3])
    drawing.add(line)

    return drawing


def _build_quality_bar_chart(rows, threshold, max_items=15):
    """Construit un HorizontalBarChart pour la qualité."""
    if not rows:
        return None

    import textwrap

                                                                                 
    rows_sorted = sorted(rows, key=lambda r: (r.get("rate") is None, r.get("rate") or 0))
    items = rows_sorted[:max_items]

    def _wrap_label(r):
                                                                                 
                                                                                 
        full = f"{r['type']} - {r['rule_name']}"
        parts = textwrap.wrap(full, width=42)
        if len(parts) > 2:
            parts = parts[:2]
            parts[1] = parts[1][:39].rstrip() + "…"
        return "\n".join(parts) if parts else full

    labels = [_wrap_label(r) for r in items]
    values = [r["rate"] if r["rate"] is not None else 0 for r in items]
    n = len(items)

    chart_height = max(140, n * 24 + 40)
    chart_width = 480

    drawing = Drawing(chart_width, chart_height)
    bc = HorizontalBarChart()
    bc.x = 205                                                             
    bc.y = 20
    bc.width = 210
    bc.height = chart_height - 40
    bc.data = [values]
    bc.categoryAxis.categoryNames = labels
    bc.categoryAxis.labels.fontSize = 6.5
    bc.categoryAxis.labels.fontName = "Helvetica"
    bc.categoryAxis.labels.dx = -4
    bc.categoryAxis.labels.textAnchor = "end"
    bc.valueAxis.valueMin = 0
    bc.valueAxis.valueMax = 100
    bc.valueAxis.valueStep = 20
    bc.valueAxis.labels.fontSize = 7
    bc.valueAxis.labels.fontName = "Helvetica"
    bc.bars[0].strokeColor = None

                                          
    bc.barLabelFormat = '%d%%'
    bc.barLabels.fontName = "Helvetica-Bold"
    bc.barLabels.fontSize = 6
    bc.barLabels.fillColor = BOA_BLUE
    bc.barLabels.boxAnchor = 'w'
    bc.barLabels.textAnchor = 'start'
    bc.barLabels.dx = 4
    bc.barLabels.nudge = 4

                                                          
    for i, rate in enumerate(values):
        bc.bars[0, i].fillColor = _rate_color(rate, threshold)

    drawing.add(bc)

    seuil_x = bc.x + (threshold / 100) * bc.width
    line = Line(seuil_x, bc.y, seuil_x, bc.y + bc.height,
                strokeColor=BOA_RED, strokeWidth=1, strokeDashArray=[3, 3])
    drawing.add(line)

    return drawing


                                                                                 
                             
                                                                                 

class _HeaderFooterCanvas(pdfcanvas.Canvas):
    """
    Canvas personnalisé qui ajoute sur chaque page :
    - En-tête bleu BOA avec titre de section (haut à droite) + logo (haut à droite)
    - Pied de page marine avec numérotation des pages
    """

    def __init__(self, *args, scope_label="", threshold=90, date_str="",
                 page_titles=None, **kwargs):
        super().__init__(*args, **kwargs)
        self._scope_label = scope_label
        self._threshold = threshold
        self._date_str = date_str
        self._current_section_title = ""
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self._draw_header_footer(num_pages)
            pdfcanvas.Canvas.showPage(self)
        pdfcanvas.Canvas.save(self)

    def _draw_header_footer(self, num_pages):
        w, h = A4
        pn = self._pageNumber                          

                                                     
                                         
                                                     
        if pn > 1:
            self.saveState()
            HEADER_H = 1.3 * cm

                                                           
            self.setFillColor(BOA_GREEN)
            self.rect(0, h - HEADER_H, w * 0.55, HEADER_H, fill=1, stroke=0)

                                                                     
            self.setFillColor(BOA_BLUE)
            self.rect(w * 0.55, h - HEADER_H, w * 0.45, HEADER_H, fill=1, stroke=0)

                                                                 
            logo = _logo_path()
            if logo:
                try:
                    logo_h = HEADER_H - 0.15 * cm
                    self.drawImage(
                        logo,
                        w - 3.2 * cm, h - HEADER_H + 0.08 * cm,
                        width=2.8 * cm,
                        height=logo_h,
                        preserveAspectRatio=True,
                        mask="auto",
                    )
                except Exception:
                    pass
            else:
                                               
                self.setFillColor(BOA_WHITE)
                self.setFont("Helvetica-Bold", 7)
                self.drawRightString(w - 0.4 * cm, h - HEADER_H + 0.55 * cm, "BOA GROUP")

                                                       
            self.setFillColor(BOA_WHITE)
            self.setFont("Helvetica-Bold", 8)
            self.drawString(0.5 * cm, h - HEADER_H + 0.70 * cm,
                            "RAPPORT DE PILOTAGE KYC")
            self.setFont("Helvetica", 7)
            self.drawString(0.5 * cm, h - HEADER_H + 0.35 * cm,
                            f"Périmètre : {self._scope_label}  |  Seuil : {self._threshold:.1f}%")

                                                                  
            section_title = getattr(self, "_current_section_title", "")
            if section_title:
                self.setFillColor(BOA_WHITE)
                self.setFont("Helvetica-Bold", 8)
                                                                                               
                self.drawCentredString(
                    w * 0.55 + (w * 0.45 - 3.5 * cm) / 2,
                    h - HEADER_H + 0.55 * cm,
                    section_title,
                )

            self.restoreState()

                                                     
                                                     
                                                     
        if pn > 1:
            self.saveState()
            FOOTER_H = 0.9 * cm

                         
            self.setFillColor(BOA_BLUE)
            self.rect(0, 0, w, FOOTER_H, fill=1, stroke=0)

                                       
            self.setStrokeColor(BOA_GREEN)
            self.setLineWidth(1.5)
            self.line(0, FOOTER_H, w, FOOTER_H)

            self.setFillColor(BOA_WHITE)
            self.setFont("Helvetica-Bold", 7)
            self.drawString(0.5 * cm, 0.32 * cm, "BOA Group")

            self.setFont("Helvetica", 7)
            self.drawCentredString(
                w / 2, 0.32 * cm,
                f"Rapport généré automatiquement le {self._date_str}"
            )

                                                      
            self.setFont("Helvetica-Bold", 8)
            self.setFillColor(BOA_GREEN_MED)
            self.drawRightString(w - 0.5 * cm, 0.32 * cm,
                                 f"Page {pn} / {num_pages}")

            self.restoreState()


                                                                                 
             
                                                                                 

                                                                             
_PDF_PAGE_TITLES = {
    2: "SOMMAIRE",
    3: "SYNTHÈSE KPI",
    4: "DÉTAIL COMPLÉTUDE",
    5: "DÉTAIL QUALITÉ",
}


def _build_filiale_rates_table(filiale_rates, threshold, styles):
    """Tableau de synthèse filiale par filiale (rapport GROUPE)."""
    headers = ["Filiale", "Clients", "Complétude", "dont PP", "dont PM",
               "Qualité", "dont PP", "dont PM", ""]
    col_widths = [3.2 * cm, 1.7 * cm, 1.8 * cm, 1.5 * cm, 1.5 * cm,
                  1.6 * cm, 1.5 * cm, 1.5 * cm, 3.0 * cm]

    def _cell(rate):
        color = _rate_color(rate, threshold)
        return Paragraph(
            _fmt_pct(rate),
            ParagraphStyle("r", fontName="Helvetica-Bold", fontSize=7,
                           textColor=color, alignment=TA_RIGHT, leading=9),
        )

    table_data = [[Paragraph(h, styles["th"]) for h in headers]]
    for fr in filiale_rates:
        total_clients = (fr.get("total_pp") or 0) + (fr.get("total_pm") or 0)
        table_data.append([
            Paragraph(fr.get("filiale", ""), styles["td_bold"]),
            Paragraph(f"{total_clients:,}".replace(",", " "), styles["td_right"]),
            _cell(fr.get("comp_global")),
            _cell(fr.get("comp_pp")),
            _cell(fr.get("comp_pm")),
            _cell(fr.get("qual_global")),
            _cell(fr.get("qual_pp")),
            _cell(fr.get("qual_pm")),
            RateBar(fr.get("comp_global"), threshold, width=55, height=7),
        ])

    t = Table(table_data, colWidths=col_widths, repeatRows=1)
    style = TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), BOA_GREEN),
        ("TEXTCOLOR", (0, 0), (-1, 0), BOA_WHITE),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, 0), 7),
        ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 1), (-1, -1), 7),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [BOA_WHITE, BOA_GRAY]),
        ("GRID", (0, 0), (-1, -1), 0.3, BOA_BORDER),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("LEFTPADDING", (0, 0), (-1, -1), 5),
        ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ("ALIGN", (1, 0), (7, -1), "RIGHT"),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ])
    for i, fr in enumerate(filiale_rates, start=1):
        cg = fr.get("comp_global")
        if cg is not None and cg < threshold:
            style.add("BACKGROUND", (0, i), (-1, i), colors.HexColor("#fff7ed"))
    t.setStyle(style)
    return t


def export_pilotage_pdf(scope_data, summary, completeness_rows, quality_rows, filiale_rates=None,
                        notations_list=None, notation_kpis=None):
    """
    Génère un rapport PDF professionnel avec logo BOA, graphiques et tableaux.
    - Logo BOA + synthèse des indicateurs clés sur la page de garde
    - Titre de section (bandeau bleu) en haut à droite
    - Pied de page avec numérotation
    Retourne un HttpResponse avec le PDF en pièce jointe.
    """
    scope = scope_data["scope"]
    selected_filiale = scope_data.get("selected_filiale", "")
    scope_label = "GROUPE" if scope == "groupe" else selected_filiale
    threshold = summary.get("threshold", 90.0)
    date_str = timezone.localtime().strftime("%d/%m/%Y à %H:%M")
    styles = _build_pdf_styles()

    buffer = io.BytesIO()
    w, h = A4

                    
    doc = BaseDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=1.5 * cm,
        rightMargin=1.5 * cm,
        topMargin=2.2 * cm,                                         
        bottomMargin=1.8 * cm,                                      
    )

    frame_normal = Frame(
        doc.leftMargin, doc.bottomMargin,
        w - doc.leftMargin - doc.rightMargin,
        h - doc.topMargin - doc.bottomMargin,
        id="normal",
    )

                                      
    frame_cover = Frame(0, 0, w, h, id="cover", leftPadding=0, rightPadding=0,
                        topPadding=0, bottomPadding=0)

    templates = [
        PageTemplate(id="cover", frames=[frame_cover], onPage=lambda c, d: None),
        PageTemplate(id="normal", frames=[frame_normal]),
    ]
    doc.addPageTemplates(templates)

    story = []

                                                             
                         
                                                             

    class CoverPage(Flowable):
        def draw(self):
            cnv = self.canv
            pw, ph = A4

                                      
            cnv.setFillColor(BOA_GREEN)
            cnv.rect(0, ph * 0.45, pw, ph * 0.55, fill=1, stroke=0)

                                              
            cnv.setFillColor(colors.HexColor("#166534"))
            path = cnv.beginPath()
            path.moveTo(0, ph * 0.45)
            path.lineTo(pw * 0.6, ph * 0.45)
            path.lineTo(0, ph * 0.30)
            path.close()
            cnv.drawPath(path, fill=1, stroke=0)

                                                               
            logo = _logo_path()
            if logo:
                try:
                    logo_h = 2.0 * cm
                    cnv.drawImage(
                        logo, pw - 5.5 * cm, ph - 2.3 * cm,
                        width=4.5 * cm,
                        height=logo_h, preserveAspectRatio=True, mask="auto"
                    )
                except Exception:
                    cnv.setFillColor(BOA_WHITE)
                    cnv.setFont("Helvetica-Bold", 9)
                    cnv.drawRightString(pw - 0.5 * cm, ph - 1.0 * cm, "BOA GROUP")

                             
            cnv.setFillColor(BOA_WHITE)
            cnv.setFont("Helvetica-Bold", 28)
            cnv.drawString(1.5 * cm, ph * 0.67, "Rapport de Pilotage")
            cnv.setFont("Helvetica-Bold", 28)
            cnv.drawString(1.5 * cm, ph * 0.62, "KYC")

                              
            cnv.setStrokeColor(BOA_GREEN_MED)
            cnv.setLineWidth(2)
            cnv.line(1.5 * cm, ph * 0.59, 12 * cm, ph * 0.59)

                        
            cnv.setFont("Helvetica", 12)
            cnv.setFillColor(BOA_GREEN_MED)
            cnv.drawString(1.5 * cm, ph * 0.555,
                         "Analyse de la Completude et de la Qualite des Donnees")

                                                                    
            cnv.setFillColor(BOA_WHITE)
            cnv.setStrokeColor(BOA_BORDER)
            cnv.setLineWidth(1)
            cnv.roundRect(1.5 * cm, ph * 0.18, pw - 3 * cm, ph * 0.16, 8, fill=1, stroke=1)

                                                           
            cols = [
                ("Filiale", scope_label, 2.5 * cm),
                ("Seuil d'analyse", f"{threshold:.1f}%", 8.5 * cm),
                ("Date de génération", date_str, 14.2 * cm),
            ]
            for label, value, col_x in cols:
                               
                cnv.setFillColor(BOA_SLATE)
                cnv.setFont("Helvetica-Bold", 7)
                cnv.drawString(col_x, ph * 0.27, label.upper())

                                    
                cnv.setFillColor(BOA_BLUE)
                cnv.setFont("Helvetica-Bold", 11)
                cnv.drawString(col_x, ph * 0.27 - 0.5 * cm, value)



            comp_rate = summary.get("completeness_rate")
            qual_rate = summary.get("quality_rate")
            low_comp = summary.get("low_completeness_count", 0)
            low_qual = summary.get("low_quality_count", 0)

            def _kpi_fg(rate, is_count=False):
                if is_count:
                    return BOA_AMBER if rate > 0 else BOA_GREEN
                if rate is None:
                    return BOA_SLATE
                return BOA_AMBER if rate < threshold else BOA_GREEN

            kpis = [
                ("TAUX DE COMPLÉTUDE", _fmt_pct(comp_rate), comp_rate, False),
                ("TAUX DE QUALITÉ", _fmt_pct(qual_rate), qual_rate, False),
                ("CHAMPS SOUS SEUIL", str(low_comp), low_comp, True),
                ("RÈGLES SOUS SEUIL", str(low_qual), low_qual, True),
            ]

            margin_x = 1.5 * cm
            gap = 0.35 * cm
            card_w = (pw - 2 * margin_x - 3 * gap) / 4
            card_h = 2.3 * cm
            band_h = 0.55 * cm
            card_top = 0.9 * cm

            cnv.setFillColor(BOA_BLUE)
            cnv.setFont("Helvetica-Bold", 8)
            cnv.drawString(margin_x, card_top + card_h + 0.2 * cm, "SYNTHÈSE DES INDICATEURS CLÉS")

            for i, (label, val_str, raw, is_count) in enumerate(kpis):
                cx = margin_x + i * (card_w + gap)
                accent = _kpi_fg(raw, is_count)


                cnv.setFillColor(BOA_WHITE)
                cnv.setStrokeColor(BOA_BORDER)
                cnv.setLineWidth(0.75)
                cnv.roundRect(cx, card_top, card_w, card_h, 4, fill=1, stroke=1)


                cnv.setFillColor(accent)
                cnv.rect(cx, card_top + card_h - band_h, card_w, band_h, fill=1, stroke=0)

                cnv.setFillColor(BOA_WHITE)
                cnv.setFont("Helvetica-Bold", 6.5)
                cnv.drawCentredString(cx + card_w / 2,
                                      card_top + card_h - band_h / 2 - 0.07 * cm, label)

                cnv.setFillColor(accent)
                cnv.setFont("Helvetica-Bold", 18)
                cnv.drawCentredString(cx + card_w / 2,
                                      card_top + (card_h - band_h) / 2 - 0.15 * cm, val_str)

        def wrap(self, *args):
            return A4

    story.append(CoverPage())

    from reportlab.platypus import NextPageTemplate, PageBreak
    story.append(NextPageTemplate("normal"))
    story.append(PageBreak())




    story.append(SetSectionTitle("DÉTAIL COMPLÉTUDE"))
    story.append(_section_header("📋  Détail Complétude — Particuliers & Entreprises", styles))
    story.append(Spacer(1, 0.4 * cm))

    def _sort_desc(rows):
        return sorted(rows, key=lambda r: r["rate"] if r.get("rate") is not None else -1, reverse=True)

    comp_rows_pp = _sort_desc([r for r in completeness_rows if r.get("type") == "PP"])
    comp_rows_pm = _sort_desc([r for r in completeness_rows if r.get("type") == "PM"])

    story.append(Paragraph("Clients PP", styles["chart_title"]))
    if comp_rows_pp:
        story.append(_build_completeness_table(comp_rows_pp, threshold, styles, scope_label, full=True))
    else:
        story.append(Paragraph("Aucune donnée disponible.", styles["body"]))
    story.append(Spacer(1, 0.5 * cm))

    story.append(Paragraph("Clients PM", styles["chart_title"]))
    if comp_rows_pm:
        story.append(_build_completeness_table(comp_rows_pm, threshold, styles, scope_label, full=True))
    else:
        story.append(Paragraph("Aucune donnée disponible.", styles["body"]))

    story.append(PageBreak())




    story.append(SetSectionTitle("QUALITÉ — PARTICULIERS"))
    story.append(_section_header("🔍  Taux de Qualité — Clients PP", styles))
    story.append(Spacer(1, 0.4 * cm))

    qual_rows_pp = _sort_desc([r for r in quality_rows if r.get("type") == "PP"])
    if qual_rows_pp:
        story.append(_build_quality_table(qual_rows_pp, threshold, styles, full=True))
    else:
        story.append(Paragraph("Aucune règle qualité pour ce périmètre.", styles["body"]))

    story.append(PageBreak())




    story.append(SetSectionTitle("QUALITÉ — ENTREPRISES"))
    story.append(_section_header("🔍  Taux de Qualité — Clients PM", styles))
    story.append(Spacer(1, 0.4 * cm))

    qual_rows_pm = _sort_desc([r for r in quality_rows if r.get("type") == "PM"])
    if qual_rows_pm:
        story.append(_build_quality_table(qual_rows_pm, threshold, styles, full=True))
    else:
        story.append(Paragraph("Aucune règle qualité pour ce périmètre.", styles["body"]))

    story.append(PageBreak())




    story.append(SetSectionTitle("NOTATION DES AGENTS"))
    story.append(_section_header("🏅  Notation des Agents", styles))
    story.append(Spacer(1, 0.4 * cm))
    story.extend(_build_notations_pdf_section(notation_kpis or {}, notations_list or [], scope, styles))


    def _make_canvas(*args, **kwargs):
        return _HeaderFooterCanvas(
            *args,
            scope_label=scope_label,
            threshold=threshold,
            date_str=date_str,
            **kwargs,
        )

    doc.build(story, canvasmaker=_make_canvas)
    pdf_bytes = buffer.getvalue()
    buffer.close()

    scope_safe = re.sub(r"[^A-Za-z0-9_-]+", "_", scope_label)
    response = HttpResponse(pdf_bytes, content_type="application/pdf")
    date_file = timezone.localtime().strftime("%Y%m%d")
    response["Content-Disposition"] = (
        f'attachment; filename="rapport_pilotage_kyc_{scope_safe}_{date_file}.pdf"'
    )
    return response


                                                                                

import re

def _section_header(title, styles):
    """Retourne un bloc titre de section avec fond vert."""
    data = [[Paragraph(title, styles["section_title"])]]
    t = Table(data, colWidths=[18 * cm])
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), BOA_BLUE),
        ("TOPPADDING", (0, 0), (-1, -1), 10),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
        ("LEFTPADDING", (0, 0), (-1, -1), 12),
        ("RIGHTPADDING", (0, 0), (-1, -1), 12),
        ("ROUNDEDCORNERS", [6, 6, 6, 6]),
    ]))
    return t


def _rate_badge(rate, threshold, styles):
    """Badge coloré pour afficher le taux."""
    if rate is None:
        return Paragraph("—", styles["td_right"])
    color = _rate_color(rate, threshold)
    color_hex = color.hexval() if hasattr(color, 'hexval') else "#64748b"
                                            
    rate_str = _fmt_pct(rate)
    style = ParagraphStyle(
        "badge", fontName="Helvetica-Bold", fontSize=8,
        textColor=color, alignment=TA_RIGHT, leading=10,
    )
    return Paragraph(rate_str, style)


def _build_completeness_table(rows, threshold, styles, scope_label, full=False):
    """Construit le tableau de complétude."""
    headers = ["Type", "Périmètre", "Champ", "Total", "Incomplets", "Taux", ""]
    col_widths = [1.2 * cm, 2.5 * cm, 4 * cm, 1.5 * cm, 1.8 * cm, 1.5 * cm, 3.5 * cm]

    table_data = [[Paragraph(h, styles["th"]) for h in headers]]
    max_rows = 50 if full else len(rows)

    for row in rows[:max_rows]:
        rate = row.get("rate")
        color = _rate_color(rate, threshold)
        rate_str = _fmt_pct(rate)
        rate_style = ParagraphStyle(
            "r", fontName="Helvetica-Bold", fontSize=8,
            textColor=color, alignment=TA_RIGHT, leading=10,
        )
        table_data.append([
            Paragraph(row.get("type", ""), styles["td_bold"]),
            Paragraph(row.get("filiale", scope_label), styles["td"]),
            Paragraph(row.get("field_label", ""), styles["td"]),
            Paragraph(str(row.get("total_clients", 0)), styles["td_right"]),
            Paragraph(str(row.get("missing_count", 0)), styles["td_right"]),
            Paragraph(rate_str, rate_style),
            RateBar(rate, threshold, width=65, height=7),
        ])

    t = Table(table_data, colWidths=col_widths, repeatRows=1)
    style = TableStyle([
                 
        ("BACKGROUND", (0, 0), (-1, 0), BOA_GREEN),
        ("TEXTCOLOR", (0, 0), (-1, 0), BOA_WHITE),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, 0), 7),
               
        ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 1), (-1, -1), 7),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [BOA_WHITE, BOA_GRAY]),
        ("GRID", (0, 0), (-1, -1), 0.3, BOA_BORDER),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("LEFTPADDING", (0, 0), (-1, -1), 5),
        ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ("ALIGN", (3, 0), (5, -1), "RIGHT"),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ])

                                   
    for i, row in enumerate(rows[:max_rows], start=1):
        if row.get("is_below_threshold"):
            style.add("BACKGROUND", (0, i), (-1, i), colors.HexColor("#fff7ed"))

    t.setStyle(style)
    return t


def _build_quality_table(rows, threshold, styles, full=False):
    """Construit le tableau de qualité."""
    headers = ["Type", "Périmètre", "Règle", "Champ", "Total", "Anomalies", "Taux", ""]
    col_widths = [1.0 * cm, 2.0 * cm, 3.0 * cm, 2.5 * cm, 1.3 * cm, 1.6 * cm, 1.4 * cm, 3.2 * cm]

    table_data = [[Paragraph(h, styles["th"]) for h in headers]]
    max_rows = 50 if full else len(rows)

    for row in rows[:max_rows]:
        rate = row.get("rate")
        color = _rate_color(rate, threshold)
        rate_str = _fmt_pct(rate)
        rate_style = ParagraphStyle(
            "r", fontName="Helvetica-Bold", fontSize=8,
            textColor=color, alignment=TA_RIGHT, leading=10,
        )
        table_data.append([
            Paragraph(row.get("type", ""), styles["td_bold"]),
            Paragraph(row.get("scope_label", ""), styles["td"]),
            Paragraph(row.get("rule_name", "")[:30], styles["td"]),
            Paragraph(row.get("field_label", "")[:25], styles["td"]),
            Paragraph(str(row.get("total_clients", 0)), styles["td_right"]),
            Paragraph(str(row.get("fail_count", 0)), styles["td_right"]),
            Paragraph(rate_str, rate_style),
            RateBar(rate, threshold, width=60, height=7),
        ])

    t = Table(table_data, colWidths=col_widths, repeatRows=1)
    style = TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), BOA_GREEN),
        ("TEXTCOLOR", (0, 0), (-1, 0), BOA_WHITE),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, 0), 7),
        ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 1), (-1, -1), 7),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [BOA_WHITE, BOA_GRAY]),
        ("GRID", (0, 0), (-1, -1), 0.3, BOA_BORDER),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("LEFTPADDING", (0, 0), (-1, -1), 5),
        ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ("ALIGN", (4, 0), (6, -1), "RIGHT"),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ])

    for i, row in enumerate(rows[:max_rows], start=1):
        if row.get("is_below_threshold"):
            style.add("BACKGROUND", (0, i), (-1, i), colors.HexColor("#fff7ed"))

    t.setStyle(style)
    return t


def _dedupe_latest_notation_per_agent(notations_list):
    """Une seule notation par agent (la plus récente) : les graphiques/répartitions
    de notation portent sur les agents évalués, pas sur l'historique complet des
    évaluations (un même agent noté plusieurs fois ne doit compter qu'une fois).
    notations_list doit déjà être trié du plus récent au plus ancien."""
    seen = set()
    result = []
    for n in notations_list:
        agent_id = getattr(getattr(n, "agent", None), "pk", None)
        if agent_id in seen:
            continue
        seen.add(agent_id)
        result.append(n)
    return result


def _build_notations_pdf_section(notation_kpis, notations_list, scope, styles):
    """Construit les flowables de la page Notation des Agents (PDF)."""
    elements = []

    total_agents = notation_kpis.get("total_agents", 0)
    excellence = notation_kpis.get("excellence_rate", 0.0)

    kpi_data = [[
        [Paragraph("AGENTS ÉVALUÉS", styles["kpi_label"]), Paragraph(str(total_agents), styles["kpi_value"])],
        [Paragraph("TAUX D'EXCELLENCE", styles["kpi_label"]), Paragraph(f"{excellence:.1f}%", styles["kpi_value"])],
    ]]
    kpi_table = Table(kpi_data, colWidths=[9 * cm] * 2)
    kpi_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), BOA_GRAY),
        ("ROUNDEDCORNERS", [8, 8, 8, 8]),
        ("BOX", (0, 0), (-1, -1), 0, colors.white),
        ("INNERGRID", (0, 0), (-1, -1), 0.5 * mm, colors.white),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 10),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
    ]))
    elements.append(kpi_table)
    elements.append(Spacer(1, 0.6 * cm))

    grades = ["Très Bien", "Bien", "Passable", "Insuffisant"]
    latest_per_agent = _dedupe_latest_notation_per_agent(notations_list)
    overall = {g: 0 for g in grades}
    by_fil = {}
    for n in latest_per_agent:
        note = getattr(n, "note", None)
        fil = getattr(getattr(n, "agent", None), "filiale", None) or "—"
        if note in overall:
            overall[note] += 1
        d = by_fil.setdefault(fil, {g: 0 for g in grades})
        if note in d:
            d[note] += 1

    if not notations_list:
        elements.append(Paragraph("Aucune notation enregistrée pour ce périmètre.", styles["body"]))
        return elements

    total_all = sum(overall.values()) or 1
    headers = ["Note", "Nombre", "Pourcentage"]
    table_data = [[Paragraph(h, styles["th"]) for h in headers]]
    for g in grades:
        table_data.append([
            Paragraph(g, styles["td_bold"]),
            Paragraph(str(overall[g]), styles["td_right"]),
            Paragraph(f"{100.0 * overall[g] / total_all:.1f}%", styles["td_right"]),
        ])
    t = Table(table_data, colWidths=[6 * cm, 4 * cm, 4 * cm], repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), BOA_GREEN),
        ("TEXTCOLOR", (0, 0), (-1, 0), BOA_WHITE),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, 0), 8),
        ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 1), (-1, -1), 8),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [BOA_WHITE, BOA_GRAY]),
        ("GRID", (0, 0), (-1, -1), 0.3, BOA_BORDER),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    elements.append(Paragraph("Répartition par type d'évaluation", styles["chart_title"]))
    elements.append(t)

    if scope == "groupe" and by_fil:
        elements.append(Spacer(1, 0.5 * cm))
        elements.append(Paragraph("Notations par filiale", styles["chart_title"]))
        headers2 = ["Filiale"] + grades + ["Total"]
        table_data2 = [[Paragraph(h, styles["th"]) for h in headers2]]
        for fil in sorted(by_fil.keys()):
            d = by_fil[fil]
            tot = sum(d.values())
            table_data2.append(
                [Paragraph(fil, styles["td_bold"])] +
                [Paragraph(str(d[g]), styles["td_right"]) for g in grades] +
                [Paragraph(str(tot), styles["td_right"])]
            )
        colw2 = [3.5 * cm] + [2.7 * cm] * 4 + [2 * cm]
        t2 = Table(table_data2, colWidths=colw2, repeatRows=1)
        t2.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), BOA_GREEN),
            ("TEXTCOLOR", (0, 0), (-1, 0), BOA_WHITE),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, 0), 7.5),
            ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
            ("FONTSIZE", (0, 1), (-1, -1), 7.5),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [BOA_WHITE, BOA_GRAY]),
            ("GRID", (0, 0), (-1, -1), 0.3, BOA_BORDER),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]))
        elements.append(t2)

    return elements






def export_pilotage_pptx(scope_data, summary, completeness_rows, quality_rows,
                         notations_list=None, notation_kpis=None, filiale_rates=None):
    """
    Génère une présentation PowerPoint reproduisant fidèlement le thème
    "Conformité BOA Group" (Calibri, vert #009A56, marine #1B2A4A) :
    couverture verte avec cercles décoratifs, en-têtes de section pleins,
    cartes KPI, jauges en anneau, graphiques natifs et tableaux stylés.
    """
    from pptx.util import Inches, Pt, Emu

    scope       = scope_data["scope"]
    fil         = scope_data.get("selected_filiale", "")
    scope_label = "GROUPE" if scope == "groupe" else (fil or "GROUPE")
    threshold   = summary.get("threshold", 90.0)
    _now        = timezone.localtime()
    date_str    = _now.strftime("%d/%m/%Y")
    _MOIS_FR    = ["", "Janvier", "Février", "Mars", "Avril", "Mai", "Juin",
                   "Juillet", "Août", "Septembre", "Octobre", "Novembre", "Décembre"]
    mois_str    = f"{_MOIS_FR[_now.month]} {_now.year}"
    _derniere_maj = summary.get("derniere_maj_kyc")
    maj_str     = f"Données KYC mises à jour le {_derniere_maj.strftime('%d/%m/%Y')}" if _derniere_maj else ""

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

                                                                                
    comp_pp = [r for r in completeness_rows if r.get("type") == "PP"]
    comp_pm = [r for r in completeness_rows if r.get("type") == "PM"]
    qual_pp = [r for r in quality_rows if r.get("type") == "PP"]
    qual_pm = [r for r in quality_rows if r.get("type") == "PM"]

    has_notation = bool(notation_kpis and notation_kpis.get("total_notations"))

                                                                                
    toc = [
        ("I",   "Complétude des données"),
        ("II",  "Qualité des données"),
    ]
    if has_notation:
        toc.append(("III", "Notation des agents"))


    _bx_cover(prs, scope_label, mois_str, maj_str)

                                                                                
    _bx_sommaire(prs, toc)

                                                            
    s = _bx_section(prs, "I", "COMPLÉTUDE DES DONNÉES")
    _bx_kpi_triple(
        s, "COMPLÉTUDE DES DONNÉES",
        ("Global", summary.get("completeness_rate"), True,
         f"{_fmt_int(summary.get('completeness_total', 0))} clients analysés · "
         f"{summary.get('low_completeness_count', 0)} champ(s) sous seuil"),
        ("PP — Particuliers", summary.get("completeness_rate_pp"), True,
         f"{_fmt_int(summary.get('completeness_total_pp', 0))} clients analysés · "
         f"{summary.get('low_completeness_count_pp', 0)} champ(s) sous seuil"),
        ("PM — Entreprises", summary.get("completeness_rate_pm"), True,
         f"{_fmt_int(summary.get('completeness_total_pm', 0))} clients analysés · "
         f"{summary.get('low_completeness_count_pm', 0)} champ(s) sous seuil"),
        threshold,
    )
    _bx_dual_charts(
        s, threshold,
        comp_pp, comp_pm, "rate", "field_label",
        "Taux par champ — Clients PP", "Taux par champ — Clients PM",
        label_max=22,
    )

                                                             
    if scope == "groupe" and filiale_rates:
        s = _bx_section(prs, "I", "COMPLÉTUDE — TAUX GLOBAL PAR FILIALE")
        _bx_sublabel(s, "TAUX DE COMPLÉTUDE PAR FILIALE (GLOBAL · PP · PM)")
        _bx_filiale_rates_chart(s, filiale_rates, "comp_global", "comp_pp", "comp_pm", threshold)




    s = _bx_section(prs, "II", "QUALITÉ DES DONNÉES")
    _bx_kpi_triple(
        s, "QUALITÉ DES DONNÉES",
        ("Global", summary.get("quality_rate"), True,
         f"{_fmt_int(summary.get('completeness_total', 0))} clients analysés · "
         f"{summary.get('low_quality_count', 0)} règle(s) sous seuil"),
        ("PP — Particuliers", summary.get("quality_rate_pp"), True,
         f"{_fmt_int(summary.get('completeness_total_pp', 0))} clients analysés · "
         f"{summary.get('low_quality_count_pp', 0)} règle(s) sous seuil"),
        ("PM — Entreprises", summary.get("quality_rate_pm"), True,
         f"{_fmt_int(summary.get('completeness_total_pm', 0))} clients analysés · "
         f"{summary.get('low_quality_count_pm', 0)} règle(s) sous seuil"),
        threshold,
    )
    _bx_dual_charts(
        s, threshold,
        qual_pp, qual_pm, "rate", "rule_name",
        "Conformité par règle — Clients PP",
        "Conformité par règle — Clients PM",
        label_max=500, label_wrap=30,
    )

                                                           
    if scope == "groupe" and filiale_rates:
        s = _bx_section(prs, "II", "QUALITÉ — TAUX GLOBAL PAR FILIALE")
        _bx_sublabel(s, "TAUX DE QUALITÉ PAR FILIALE (GLOBAL · PP · PM)")
        _bx_filiale_rates_chart(s, filiale_rates, "qual_global", "qual_pp", "qual_pm", threshold)

                                                            
    if has_notation:
        s = _bx_section(prs, "III", "NOTATION DES AGENTS")
        _bx_notation(s, notation_kpis, notations_list or [], scope=scope)

                                                                                
    _bx_end(prs, scope_label, mois_str)

                                                                                
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    pptx_bytes = buf.getvalue()
    buf.close()

    scope_safe = re.sub(r"[^A-Za-z0-9_-]+", "_", scope_label)
    date_file  = timezone.localtime().strftime("%Y%m%d")
    response = HttpResponse(
        pptx_bytes,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    response["Content-Disposition"] = (
        f'attachment; filename="rapport_pilotage_kyc_{scope_safe}_{date_file}.pptx"'
    )
    return response


                                                                                 
                                   
                                                                                 

BX_GREEN      = RGBColor(0x00, 0x9A, 0x56)                       
BX_GREEN_BRT  = RGBColor(0x00, 0xC0, 0x60)                                  
BX_GREEN_PALE = RGBColor(0xC8, 0xE6, 0xC9)                        
BX_GREEN_SUB  = RGBColor(0xC8, 0xF0, 0xD8)                          
BX_NAVY       = RGBColor(0x1B, 0x2A, 0x4A)               
BX_BODY       = RGBColor(0x33, 0x33, 0x33)                   
BX_MUTED      = RGBColor(0x66, 0x66, 0x66)                     
BX_FOOT       = RGBColor(0x88, 0x88, 0x88)                 
BX_CARD       = RGBColor(0xF5, 0xF5, 0xF5)               
BX_DIVIDER    = RGBColor(0xE0, 0xE0, 0xE0)           
BX_TRACK      = RGBColor(0xEC, 0xEF, 0xF1)                          
BX_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BX_RED        = RGBColor(0xD3, 0x2F, 0x2F)                 
BX_AMBER      = RGBColor(0xF5, 0xA6, 0x23)                    
BX_CYAN       = RGBColor(0x29, 0xAB, 0xE2)                

_FONT = "Calibri"
_FOOTER_TXT = "Conformité BOA Group"


def _roman(n):
    return ["", "I", "II", "III", "IV", "V", "VI", "VII", "VIII", "IX", "X"][n] if n < 11 else str(n)


def _trunc(s, n):
    s = str(s or "")
    return s if len(s) <= n else s[: n - 1] + "…"


def _fmt_int(v):
    try:
        return f"{int(v):,}".replace(",", " ")
    except (ValueError, TypeError):
        return str(v)


def _bx_rate_color(rate, threshold):
    """Vert si le taux atteint/dépasse le seuil, orange en deçà. Pas d'autre couleur."""
    if rate is None:
        return BX_MUTED
    return BX_GREEN if rate >= threshold else BX_AMBER


def _bx_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bx_rect(slide, l, t, w, h, fill, line=None, line_w=1.0):
    """Rectangle plein, sans ombre (Inches)."""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is not None:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _bx_oval(slide, l, t, w, h, fill):
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _bx_ring(slide, l, t, w, h, line_color, line_w=1.5):
    """Cercle décoratif simple : contour seul, sans remplissage."""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.background()
    sp.line.color.rgb = line_color
    sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def _bx_text(slide, l, t, w, h, text, size, color, bold=False, align=None,
             anchor=None, italic=False, wrap=True, font=_FONT, spacing=None):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_top = tf.margin_bottom = 0
    tf.margin_left = tf.margin_right = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = font
        p.font.color.rgb = color
        if align is not None:
            p.alignment = align
        if spacing is not None:
            p.line_spacing = spacing
    return tb


def _bx_logo(slide):
    """Boîte logo marine BOA en haut à droite (calquée sur l'exemple)."""
    from pptx.util import Inches
    logo = _logo_path()
    if logo:
        try:
            slide.shapes.add_picture(logo, Inches(7.45), Inches(0.0), width=Inches(2.55))
            return
        except Exception:
            pass
                                   
    _bx_rect(slide, 7.45, 0.0, 2.55, 0.62, BX_NAVY)
    _bx_text(slide, 7.55, 0.06, 2.4, 0.5, "BANK OF AFRICA\nBMCE GROUP", 9, BX_WHITE,
             bold=True, align=None)


def _bx_decor(slide, color):
    """Motif décoratif épuré : deux anneaux fins en haut à droite."""
    _bx_ring(slide, 8.7, -1.1, 3.6, 3.6, color, 1.5)
    _bx_ring(slide, 9.3, 1.4, 2.2, 2.2, color, 1.5)


def _bx_footer(slide):
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    _bx_rect(slide, 0.0, 5.35, 10.0, 0.005, BX_DIVIDER)
    _bx_text(slide, 6.0, 5.37, 3.85, 0.25, _FOOTER_TXT, 8, BX_FOOT,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _bx_section(prs, roman, title):
    """En-tête de section : bande verte pleine + cercles + logo + pied."""
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    s = _bx_slide(prs)
    _bx_decor(s, BX_GREEN_PALE)
    _bx_rect(s, 0.0, 0.0, 10.0, 0.62, BX_GREEN)
    _bx_text(s, 0.4, 0.0, 7.0, 0.62, f"{roman}.  {title}", 18, BX_WHITE,
             bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _bx_logo(s)
    _bx_footer(s)
    return s


def _bx_sublabel(slide, text, top=0.82):
    """Sous-titre vert majuscule + filet vert pleine largeur."""
    from pptx.enum.text import MSO_ANCHOR
    _bx_text(slide, 0.4, top, 9.2, 0.3, text, 11, BX_GREEN, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    _bx_rect(slide, 0.4, top + 0.3, 9.2, 0.035, BX_GREEN)


                                                                                 
                               
                                                                                 

def _bx_cover(prs, scope_label, mois_str, maj_str=""):
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    s = _bx_slide(prs)
    _bx_rect(s, 0.0, 0.0, 10.0, 5.625, BX_GREEN)
    _bx_ring(s, 7.9, -1.6, 4.6, 4.6, BX_GREEN_BRT, 2.0)
    _bx_ring(s, 8.9, 2.0, 2.8, 2.8, BX_GREEN_BRT, 2.0)
    _bx_logo(s)

    _bx_text(s, 0.5, 1.35, 7.5, 1.5, "RAPPORT DE PILOTAGE KYC", 46, BX_WHITE,
             bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _bx_text(s, 0.5, 2.9, 7.5, 0.5,
             f"SYNTHÈSE QUALITÉ & COMPLÉTUDE — {scope_label.upper()}",
             14, BX_GREEN_SUB, bold=True)
    _bx_rect(s, 0.52, 3.45, 6.0, 0.045, BX_WHITE)
    _bx_text(s, 0.5, 3.65, 7.0, 0.35, mois_str, 11, RGBColor(0xD4, 0xF5, 0xE4))
    if maj_str:
        _bx_text(s, 0.5, 3.98, 7.5, 0.3, maj_str, 10, RGBColor(0xD4, 0xF5, 0xE4), italic=True)
    _bx_text(s, 0.5, 5.2, 7.0, 0.28, _FOOTER_TXT, 8,
             RGBColor(0xA0, 0xD4, 0xB8), italic=True)


def _bx_sommaire(prs, toc):
    from pptx.enum.text import MSO_ANCHOR
    s = _bx_slide(prs)
    _bx_decor(s, BX_GREEN_PALE)
    _bx_logo(s)
    _bx_footer(s)
    _bx_text(s, 0.5, 0.5, 6.0, 1.1, "SOMMAIRE", 52, BX_GREEN, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

    n = len(toc)
    top = 1.85
    row_h = min(0.62, (4.9 - top) / max(n, 1))
    for roman, title in toc:
        _bx_rect(s, 0.5, top + 0.04, 0.07, 0.38, BX_GREEN)
        _bx_text(s, 0.68, top, 0.6, 0.42, roman, 12, BX_GREEN, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        _bx_text(s, 1.3, top, 7.0, 0.42, title, 14, BX_NAVY, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        _bx_rect(s, 0.5, top + row_h - 0.12, 7.6, 0.02, BX_DIVIDER)
        top += row_h


def _bx_kpi_card(slide, l, t, w, h, label, value, color, sub=None):
    """Carte KPI : bandeau coloré + corps gris clair + grand chiffre."""
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    _bx_rect(slide, l, t, w, h, BX_CARD)
    _bx_rect(slide, l, t, w, 0.34, color)
    _bx_text(slide, l + 0.05, t, w - 0.1, 0.34, label, 8.5, BX_WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    val_h = (h - 0.34) if not sub else (h - 0.34 - 0.3)
    _bx_text(slide, l, t + 0.34, w, val_h, value, 33, color, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        _bx_text(slide, l, t + h - 0.32, w, 0.28, sub, 8.5, BX_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _bx_donut(slide, l, t, sz, rate, threshold, title):
    """Jauge en anneau native + grand pourcentage centré + titre dessous."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    color = _bx_rate_color(rate, threshold)
    r = max(0.0, min(float(rate or 0), 100.0))
    cd = CategoryChartData()
    cd.categories = ["ok", "reste"]
    cd.add_series("s", (r, 100.0 - r))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(l), Inches(t),
                                Inches(sz), Inches(sz), cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.has_data_labels = False
    try:
        from pptx.oxml.ns import qn
        plot._element.find(qn("c:firstSliceAng"))
        hole = plot._element.find(qn("c:holeSize"))
        if hole is not None:
            hole.set("val", "68")
    except Exception:
        pass
    pts = plot.series[0].points
    pts[0].format.fill.solid()
    pts[0].format.fill.fore_color.rgb = color
    pts[1].format.fill.solid()
    pts[1].format.fill.fore_color.rgb = BX_TRACK
    for pt in pts:
        pt.format.line.fill.background()

                                              
    _bx_text(slide, l, t + sz / 2 - 0.32, sz, 0.6,
             _fmt_pct(r), 26, color, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _bx_text(slide, l - 0.4, t + sz + 0.05, sz + 0.8, 0.3, title, 11, BX_NAVY,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)




def _bx_hbar_chart(slide, rows, threshold, value_key, label_key, type_key=None,
                   top=1.55, left=0.4, width=9.2, max_items=11, label_max=34,
                   font_sz=8.5, max_h=3.55, vmax=100, return_bottom=False, label_wrap=None):
    """Graphique à barres horizontales natif, points colorés selon le seuil.
    label_wrap : si défini, affiche le libellé entier sur plusieurs lignes (largeur en
    caractères) au lieu de le tronquer avec label_max."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

    items = list(rows)[:max_items]
    chart_h = min(max_h, 0.30 * max(len(items), 1) + 0.5)
    if not items:
        _bx_text(slide, left, top + chart_h / 2 - 0.2, width, 0.4,
                 "Aucune donnée disponible.", 11, BX_MUTED, align=PP_ALIGN.CENTER)
        return top + chart_h if return_bottom else None

    cats, vals, colors = [], [], []
    for it in items:
        rate = it.get(value_key) or 0
        typ = str(it.get(type_key, "")).strip() if type_key else ""
        raw_label = str(it.get(label_key, "") or "")
        if label_wrap:
            import textwrap
            lab = "\n".join(textwrap.wrap(raw_label, width=label_wrap)) or raw_label
        else:
            lab = _trunc(raw_label, label_max)
        cats.append(f"{typ}  {lab}" if typ else lab)
        vals.append(round(float(rate), 1))
        colors.append(_bx_rate_color(rate, threshold))

    cats.reverse(); vals.reverse(); colors.reverse()               

    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("Taux", vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                                Inches(left), Inches(top),
                                Inches(width), Inches(chart_h), cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False

    plot = chart.plots[0]
    plot.gap_width = 55
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = '0.#"%"'
    dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.size = Pt(font_sz)
    dl.font.bold = True
    dl.font.name = _FONT
    dl.font.color.rgb = BX_BODY

    series = plot.series[0]
    for idx, pt in enumerate(series.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[idx]
        pt.format.line.fill.background()

    va = chart.value_axis
    va.minimum_scale = 0
    va.maximum_scale = vmax
    va.has_major_gridlines = False
    va.visible = False
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(font_sz)
    ca.tick_labels.font.name = _FONT
    ca.tick_labels.font.color.rgb = BX_BODY
    ca.format.line.color.rgb = BX_DIVIDER

    return top + chart_h if return_bottom else None


def _bx_legend(slide, y, threshold, x=0.4):
    from pptx.enum.text import MSO_ANCHOR
    items = [(BX_GREEN, f"Conforme (≥ {threshold:.0f}%)"),
             (BX_AMBER, f"Sous le seuil (< {threshold:.0f}%)")]
    for col, lab in items:
        _bx_rect(slide, x, y + 0.04, 0.2, 0.2, col)
        _bx_text(slide, x + 0.28, y, 2.4, 0.3, lab, 9, BX_MUTED,
                 anchor=MSO_ANCHOR.MIDDLE)
        x += 2.7


def _bx_legend_items(slide, y, entries, x=0.4, font_sz=9, step=1.7):
    """Légende générique : entries = [(color, label), ...]."""
    from pptx.enum.text import MSO_ANCHOR
    for col, lab in entries:
        _bx_rect(slide, x, y + 0.04, 0.2, 0.2, col)
        _bx_text(slide, x + 0.28, y, step - 0.3, 0.3, lab, font_sz, BX_MUTED,
                 anchor=MSO_ANCHOR.MIDDLE)
        x += step


def _bx_table(slide, headers, rows_data, col_fracs, source_rows, threshold, rate_key,
              top=1.5):
    """Tableau stylé : en-tête vert, alternance de lignes, taux coloré."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    n_rows = len(rows_data) + 1
    n_cols = len(headers)
    total_w = 9.2
    left = 0.4
    col_w = [f * total_w for f in col_fracs]
    head_h = 0.42
    row_h = min(0.4, (3.6 - head_h) / max(len(rows_data), 1))
    height = head_h + row_h * len(rows_data)

    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                 Inches(total_w), Inches(height)).table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, w in enumerate(col_w):
        tbl.columns[i].width = Inches(w)
    tbl.rows[0].height = Inches(head_h)
    for i in range(1, n_rows):
        tbl.rows[i].height = Inches(row_h)

    right_cols = {n_cols - 1, n_cols - 2}
             
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_top = c.margin_bottom = Inches(0.02)
        c.margin_left = c.margin_right = Inches(0.07)
        c.fill.solid(); c.fill.fore_color.rgb = BX_GREEN
        p = c.text_frame.paragraphs[0]
        p.font.bold = True; p.font.size = Pt(9); p.font.name = _FONT
        p.font.color.rgb = BX_WHITE
        p.alignment = PP_ALIGN.CENTER if j == 0 else (PP_ALIGN.RIGHT if j in right_cols else PP_ALIGN.LEFT)

    for i, vals in enumerate(rows_data, start=1):
        src = source_rows[i - 1] if i - 1 < len(source_rows) else {}
        bg = BX_WHITE if i % 2 else RGBColor(0xF2, 0xF6, 0xFA)
        for j, val in enumerate(vals):
            c = tbl.cell(i, j)
            c.text = str(val)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = c.margin_bottom = Inches(0.01)
            c.margin_left = c.margin_right = Inches(0.07)
            c.fill.solid(); c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(8.5); p.font.name = _FONT
            if j == n_cols - 1 and rate_key is not None:
                p.font.bold = True
                p.font.color.rgb = _bx_rate_color(src.get(rate_key), threshold)
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.font.color.rgb = BX_BODY
                p.alignment = PP_ALIGN.CENTER if j == 0 else (PP_ALIGN.RIGHT if j in right_cols else PP_ALIGN.LEFT)


def _bx_end(prs, scope_label, mois_str):
    from pptx.enum.text import MSO_ANCHOR
    s = _bx_slide(prs)
    _bx_rect(s, 0.0, 0.0, 10.0, 5.625, BX_GREEN)
    _bx_ring(s, 7.9, -1.6, 4.6, 4.6, BX_GREEN_BRT, 2.0)
    _bx_ring(s, 8.9, 2.0, 2.8, 2.8, BX_GREEN_BRT, 2.0)
    _bx_logo(s)
    _bx_text(s, 0.5, 1.7, 7.5, 1.2, "FIN DU DOCUMENT", 44, BX_WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    _bx_rect(s, 0.52, 3.05, 6.0, 0.045, BX_WHITE)
    _bx_text(s, 0.5, 3.25, 7.5, 0.4,
             f"RAPPORT DE PILOTAGE KYC — {_FOOTER_TXT}", 13, BX_WHITE, bold=True)
    _bx_text(s, 0.5, 3.85, 7.0, 0.35, mois_str, 11,
             RGBColor(0xD4, 0xF5, 0xE4), italic=True)


def _bx_kpi_triple(slide, header, c1, c2, c3, threshold):
    """Carte KPI 3 colonnes (Global / PP / PM) calquée sur la page pilotage.
    Chaque colonne ci = (label, valeur_float_ou_None, est_pourcentage, sous_texte_optionnel)
    — le sous-texte (ex. "12 345 clients · 3 sous seuil") s'affiche sous chaque bloc."""
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    L, T, W = 0.4, 1.1, 9.2
    HEAD, BODY, FOOT = 0.32, 0.82, 0.24
    H = HEAD + BODY + FOOT

    _bx_rect(slide, L, T, W, H, BX_WHITE, line=BX_DIVIDER, line_w=0.75)
    _bx_rect(slide, L, T, W, HEAD, BX_NAVY)
    _bx_oval(slide, L + 0.16, T + HEAD / 2 - 0.06, 0.12, 0.12, BX_GREEN)
    _bx_text(slide, L + 0.4, T, W - 0.5, HEAD, header.upper(), 9.5, BX_GREEN_SUB,
             bold=True, anchor=MSO_ANCHOR.MIDDLE)

    cols = [c1, c2, c3]
    cw = W / 3.0
    fy = T + HEAD + BODY + 0.02
    for i, col in enumerate(cols):
        label, val, is_pct = col[0], col[1], col[2]
        foot = col[3] if len(col) > 3 else ""
        cx = L + i * cw
        if i > 0:
            _bx_rect(slide, cx, T + HEAD + 0.08, 0.008, BODY - 0.16, BX_DIVIDER)
        _bx_text(slide, cx, T + HEAD + 0.06, cw, 0.2,
                 label.upper(), 8, BX_MUTED, bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        if val is None:
            valtxt, color = "—", BX_MUTED
        elif is_pct:
            valtxt = _fmt_pct(val)
            color = _bx_rate_color(val, threshold)
        else:
            valtxt, color = _fmt_int(val), BX_NAVY
        big = 30 if i == 0 else 26
        _bx_text(slide, cx, T + HEAD + 0.24, cw, 0.4, valtxt, big, color,
                 bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if is_pct and val is not None:
            bw = 1.2
            bx0 = cx + (cw - bw) / 2
            by0 = T + HEAD + 0.66
            _bx_rect(slide, bx0, by0, bw, 0.08, BX_TRACK)
            _bx_rect(slide, bx0, by0, bw * max(0, min(val, 100)) / 100.0, 0.08, color)
        if foot:
            _bx_text(slide, cx + 0.05, fy, cw - 0.1, FOOT, foot, 7.5, BX_MUTED,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _bx_dual_charts(slide, threshold, rows_l, rows_r, value_key, label_key,
                    title_l, title_r, label_max=22, max_items=None, label_wrap=None):
    """Deux graphiques à barres côte à côte (PP / PM) + légende de seuil partagée.
    max_items=None affiche tous les champs (dimensionnement dynamique).
    label_wrap : voir _bx_hbar_chart (affiche les libellés entiers sur plusieurs lignes)."""
    from pptx.enum.text import MSO_ANCHOR

    def _worst_first(rows):
        return sorted(rows, key=lambda r: (r.get(value_key) is None, r.get(value_key) or 0))

    n = max(len(rows_l), len(rows_r), 1)
    cap = n if max_items is None else min(n, max_items)
    fs = 7.0 if cap <= 8 else (6.5 if cap <= 12 else 6.0)
    ctop, bottom = 2.82, 4.92
    max_h = bottom - ctop

    for left, rows, title in ((0.4, rows_l, title_l), (5.15, rows_r, title_r)):
        _bx_text(slide, left, 2.52, 4.45, 0.26, title, 9.5, BX_NAVY, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        _bx_hbar_chart(slide, _worst_first(rows), threshold, value_key, label_key,
                       type_key=None, top=ctop, left=left, width=4.45,
                       max_items=cap, label_max=label_max, font_sz=fs, max_h=max_h, vmax=120,
                       label_wrap=label_wrap)

    _bx_legend(slide, 4.98, threshold, x=2.6)


def _bx_notation(slide, kpis, notations_list, scope="groupe"):
    """Slide notation : 2 KPI + (groupe) barres par filiale + anneau de répartition,
    (filiale) uniquement l'anneau de répartition par type d'évaluation."""
    from pptx.enum.text import MSO_ANCHOR

    _bx_sublabel(slide, "PERFORMANCE DES AGENTS NOTÉS")

    total_agents = kpis.get("total_agents", 0)
    excellence   = kpis.get("excellence_rate", 0.0)

    cw, gap = 4.53, 0.145
    _bx_kpi_card(slide, 0.4, 1.2, cw, 1.15, "AGENTS ÉVALUÉS", _fmt_int(total_agents),
                 BX_NAVY, sub="chargé(s) de compte")
    _bx_kpi_card(slide, 0.4 + (cw + gap), 1.2, cw, 1.15, "TAUX D'EXCELLENCE",
                 f"{excellence:.1f}%", BX_GREEN, sub="notes « Bien » & « Très Bien »")

    grades = ["Très Bien", "Bien", "Passable", "Insuffisant"]
    gcolor = {"Très Bien": RGBColor(0x10, 0xB9, 0x81),
              "Bien": RGBColor(0x3B, 0x82, 0xF6),
              "Passable": RGBColor(0xF5, 0x9E, 0x0B),
              "Insuffisant": RGBColor(0xEF, 0x44, 0x44)}
    latest_per_agent = _dedupe_latest_notation_per_agent(notations_list)
    overall = {g: 0 for g in grades}
    by_fil = {}
    for n in latest_per_agent:
        note = getattr(n, "note", None)
        fil = getattr(getattr(n, "agent", None), "filiale", None) or "—"
        if note in overall:
            overall[note] += 1
        d = by_fil.setdefault(fil, {g: 0 for g in grades})
        if note in d:
            d[note] += 1

    labels = [g for g in grades if overall[g] > 0]
    values = [overall[g] for g in labels]
    colors = [gcolor[g] for g in labels]

    if scope == "filiale":
                                                                               
        _bx_text(slide, 0.4, 2.6, 9.2, 0.3, "RÉPARTITION PAR TYPE D'ÉVALUATION",
                 11, BX_GREEN, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        if values:
            _bx_donut_multi(slide, 3.0, 3.05, 1.95, labels, values, colors)
            ly = 3.2
            for g in grades:
                if overall[g] == 0:
                    continue
                pct = 100.0 * overall[g] / max(sum(values), 1)
                _bx_rect(slide, 5.6, ly + 0.03, 0.2, 0.2, gcolor[g])
                _bx_text(slide, 5.9, ly, 3.0, 0.28,
                         f"{g}  ·  {overall[g]}  ({pct:.0f}%)", 11, BX_BODY,
                         anchor=MSO_ANCHOR.MIDDLE)
                ly += 0.42
        else:
            _bx_text(slide, 0.4, 3.6, 9.2, 0.4, "Aucune notation.", 12, BX_MUTED)
        return

                                                                    
    _bx_text(slide, 0.4, 2.55, 4.45, 0.28, "Notations par filiale", 10, BX_NAVY,
             bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _bx_text(slide, 5.15, 2.55, 4.45, 0.28, "Répartition par type d'évaluation",
             10, BX_NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    fils = sorted(by_fil.keys())[:6]
    series = [(g, gcolor[g], [by_fil[f][g] for f in fils]) for g in grades]
    _bx_stacked_col(slide, 0.4, 2.9, 4.45, 1.95, [_trunc(f, 10) for f in fils],
                    series, data_labels=True)

    if values:
        _bx_donut_multi(slide, 5.35, 2.85, 1.8, labels, values, colors)
        ly = 2.95
        for g in grades:
            if overall[g] == 0:
                continue
            _bx_rect(slide, 7.45, ly + 0.03, 0.18, 0.18, gcolor[g])
            _bx_text(slide, 7.7, ly, 2.0, 0.26, f"{g}  ·  {overall[g]}", 9, BX_BODY,
                     anchor=MSO_ANCHOR.MIDDLE)
            ly += 0.34
    else:
        _bx_text(slide, 5.15, 3.6, 4.45, 0.4, "Aucune notation.", 11, BX_MUTED)


def _bx_stacked_col(slide, l, t, w, h, categories, series, data_labels=False):
    """Histogramme à colonnes empilées natif. series = [(nom, couleur, [valeurs]), ...]."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    if not categories:
        _bx_text(slide, l, t + h / 2 - 0.2, w, 0.4, "Aucune notation.", 11, BX_MUTED,
                 align=PP_ALIGN.CENTER)
        return

    cd = CategoryChartData()
    cd.categories = categories
    for name, _color, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED,
                                Inches(l), Inches(t), Inches(w), Inches(h), cd)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 60
    for idx, (_name, color, _vals) in enumerate(series):
        sr = plot.series[idx]
        sr.format.fill.solid()
        sr.format.fill.fore_color.rgb = color
        sr.format.line.fill.background()
    if data_labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = '0;-0;'                              
        dl.number_format_is_linked = False
        dl.font.size = Pt(8)
        dl.font.bold = True
        dl.font.name = _FONT
        dl.font.color.rgb = BX_WHITE
    va = chart.value_axis
    va.has_major_gridlines = False
    va.visible = False
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(8)
    ca.tick_labels.font.name = _FONT
    ca.tick_labels.font.color.rgb = BX_BODY
    ca.format.line.color.rgb = BX_DIVIDER


def _bx_donut_multi(slide, l, t, sz, labels, values, colors):
    """Anneau multi-catégories natif (sans étiquettes, légende externe)."""
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("Répartition", values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(l), Inches(t),
                                Inches(sz), Inches(sz), cd)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = False
    try:
        from pptx.oxml.ns import qn
        hole = plot._element.find(qn("c:holeSize"))
        if hole is not None:
            hole.set("val", "62")
    except Exception:
        pass
    pts = plot.series[0].points
    for i, pt in enumerate(pts):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        pt.format.line.fill.background()


def _bx_grouped_col(slide, l, t, w, h, categories, series, vmax=105, font_sz=8):
    """Histogramme à colonnes groupées natif (clustered). series = [(nom, couleur, [valeurs]), ...].
    Étiquettes de données en pourcentage, axe masqué."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

    if not categories:
        _bx_text(slide, l, t + h / 2 - 0.2, w, 0.4, "Aucune donnée par filiale.", 11, BX_MUTED,
                 align=PP_ALIGN.CENTER)
        return

    cd = CategoryChartData()
    cd.categories = categories
    for name, _color, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(l), Inches(t), Inches(w), Inches(h), cd)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.overlap = -10
    for idx, (_name, color, _vals) in enumerate(series):
        sr = plot.series[idx]
        sr.format.fill.solid()
        sr.format.fill.fore_color.rgb = color
        sr.format.line.fill.background()
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = '0"%";;'                                
    dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.size = Pt(font_sz)
    dl.font.bold = True
    dl.font.name = _FONT
    dl.font.color.rgb = BX_BODY
    va = chart.value_axis
    va.minimum_scale = 0
    va.maximum_scale = vmax
    va.has_major_gridlines = False
    va.visible = False
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(font_sz + 0.5)
    ca.tick_labels.font.name = _FONT
    ca.tick_labels.font.color.rgb = BX_BODY
    ca.format.line.color.rgb = BX_DIVIDER


def _bx_filiale_rates_chart(slide, rates, gkey, ppkey, pmkey, threshold):
    """Colonnes groupées Global / PP / PM par filiale + légende."""
    fils = [r.get("filiale", "—") for r in rates][:12]
    n = len(fils)

    def col(key):
        return [(r.get(key) if r.get(key) is not None else 0) for r in rates[:12]]

    series = [
        ("Global", BX_NAVY, col(gkey)),
        ("PP", BX_GREEN, col(ppkey)),
        ("PM", BX_CYAN, col(pmkey)),
    ]
    fs = 8.0 if n <= 7 else (7.0 if n <= 10 else 6.0)
    _bx_grouped_col(slide, 0.4, 1.6, 9.2, 3.25,
                    [_trunc(f, 12) for f in fils], series, vmax=108, font_sz=fs)
    _bx_legend_items(slide, 5.0,
                     [(BX_NAVY, "Global"), (BX_GREEN, "PP — Particuliers"),
                      (BX_CYAN, "PM — Entreprises")],
                     x=2.0, font_sz=9, step=2.3)