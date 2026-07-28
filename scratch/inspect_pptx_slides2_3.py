import sys
from pptx import Presentation

def inspect_range(filepath, start, end):
    prs = Presentation(filepath)
    if hasattr(sys.stdout, 'reconfigure'):
        sys.stdout.reconfigure(encoding='utf-8')
        
    for i in range(start-1, min(end, len(prs.slides))):
        slide = prs.slides[i]
        print(f"\n================ SLIDE {i+1} ================")
        for j, shape in enumerate(slide.shapes):
            shape_type = "Shape"
            if shape.shape_type == 1: shape_type = "Rectangle"
            elif shape.shape_type == 17: shape_type = "TextBox"
            elif shape.shape_type == 19: shape_type = "Table"
            elif shape.shape_type == 3: shape_type = "Chart"
            elif shape.shape_type == 13: shape_type = "Picture"
            
            text = ""
            if shape.has_text_frame:
                text = shape.text_frame.text.strip().replace('\n', ' | ')[:100]
                
            fill_color = "N/A"
            try:
                if shape.fill and shape.fill.type == 1:        
                    fill_color = str(shape.fill.fore_color.rgb)
            except Exception:
                pass
                
            print(f"[{j+1}] {shape_type} | L={shape.left.cm:.2f}cm, T={shape.top.cm:.2f}cm, W={shape.width.cm:.2f}cm, H={shape.height.cm:.2f}cm | Fill={fill_color} | Text='{text}'")

if __name__ == "__main__":
    import sys
    filepath = "scratch/test_output.pptx"
    if len(sys.argv) > 1:
        filepath = sys.argv[1]
    inspect_range(filepath, 4, 4)
